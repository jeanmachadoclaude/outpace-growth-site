"""
Outpace — Sessão Estratégica Closer
Deck de apresentação dos produtos em identidade visual do site.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- BRAND TOKENS ----------
BG        = RGBColor(0x08, 0x0F, 0x09)   # --bg
BG2       = RGBColor(0x0B, 0x1A, 0x0D)   # --bg2
BG3       = RGBColor(0x0F, 0x21, 0x15)   # --bg3
CARD      = RGBColor(0x10, 0x1E, 0x12)   # --card
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)   # --gold
GOLD_L    = RGBColor(0xE2, 0xC5, 0x6A)   # --gold-l
GOLD_D    = RGBColor(0xA0, 0x7B, 0x2A)   # --gold-d
CREAM     = RGBColor(0xED, 0xE8, 0xDB)   # --cream
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x6E, 0x88, 0x72)   # --muted
LIGHT     = RGBColor(0xB8, 0xCD, 0xB9)   # --light
BORDER    = RGBColor(0x3A, 0x36, 0x22)   # approx border

FONT_H    = "Playfair Display"   # fallback Georgia
FONT_B    = "DM Sans"            # fallback Calibri

IMG = "/Users/jeanmachadopazdasilva/sellas-website/home/images"

# ---------- SETUP ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ---------- HELPERS ----------

def add_slide():
    s = prs.slides.add_slide(BLANK)
    # base background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s

def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def oval(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def _normalize_run(r):
    """Run can be: 'str'  |  ('str',)  |  ('str', {opts})  |  'str'(with default opts)."""
    if isinstance(r, str):
        return r, {}
    if isinstance(r, tuple):
        if len(r) == 1: return r[0], {}
        if len(r) >= 2: return r[0], (r[1] or {})
    raise ValueError(f"bad run spec: {r!r}")

def _expand_paragraphs(content):
    """
    Normalize 'content' to a list of paragraphs, where each paragraph is a list of (str, opts).
    Accepts: str, list[str|list[runs]|str-with-\\n].
    Splits any '\\n' within plain strings into separate paragraphs.
    """
    if isinstance(content, str):
        content = [content]
    paragraphs = []
    for ln in content:
        if isinstance(ln, list):
            # list of runs → one paragraph (no \n splitting inside rich runs)
            paragraphs.append([_normalize_run(r) for r in ln])
        else:
            # plain string or tuple: may contain \n → multiple paragraphs
            txt, opts = _normalize_run(ln if isinstance(ln, tuple) else (ln, {}))
            parts = txt.split("\n")
            for part in parts:
                paragraphs.append([(part, opts)])
    return paragraphs

def text(slide, x, y, w, h, content, *, font=FONT_B, size=14, bold=False, italic=False,
         color=CREAM, align="left", valign="top", tracking=None, line_spacing=1.25):
    """Add text box with rich paragraph support."""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    paragraphs = _expand_paragraphs(content)
    for i, runs_spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center": p.alignment = PP_ALIGN.CENTER
        elif align == "right": p.alignment = PP_ALIGN.RIGHT
        else: p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing

        for r_i, (rt, over) in enumerate(runs_spec):
            if r_i == 0 and not p.runs:
                run = p.add_run()
            else:
                run = p.add_run()
            run.text = rt
            f = run.font
            f.name = over.get("font", font)
            f.size = Pt(over.get("size", size))
            f.bold = over.get("bold", bold)
            f.italic = over.get("italic", italic)
            f.color.rgb = over.get("color", color)
            ts = over.get("tracking", tracking)
            if ts is not None:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(ts)))
    return tx

def kicker(slide, x, y, t, color=GOLD, size=10, w=6):
    return text(slide, x, y, w, 0.3, t.upper(), size=size, bold=True,
                color=color, tracking=300, font=FONT_B)

def title_serif(slide, x, y, w, h, content, size=44, color=CREAM, align="left", line_spacing=1.05):
    return text(slide, x, y, w, h, content, font=FONT_H, size=size, bold=True,
                color=color, align=align, line_spacing=line_spacing)

def body(slide, x, y, w, h, t, *, size=13, color=LIGHT, align="left", line_spacing=1.55):
    return text(slide, x, y, w, h, t, font=FONT_B, size=size, color=color,
                align=align, line_spacing=line_spacing)

def gold_dot(slide, x, y, d=0.08):
    return oval(slide, x, y, d, d, fill=GOLD)

def gold_rule(slide, x, y, w=0.6, h=0.02):
    return rect(slide, x, y, w, h, fill=GOLD)

def footer(slide, page, total):
    text(slide, 0.6, 7.12, 3, 0.28, "OUTPACE  ·  SESSÃO ESTRATÉGICA", font=FONT_B, size=8,
         color=MUTED, tracking=250, bold=True)
    text(slide, 10, 7.12, 2.8, 0.28, f"{page:02d}  /  {total:02d}", font=FONT_B, size=8,
         color=MUTED, tracking=250, bold=True, align="right")

def header_mark(slide):
    # top-left wordmark
    text(slide, 0.6, 0.42, 2, 0.4, "Outpace", font=FONT_H, size=18, bold=True, color=GOLD)
    # hairline under nav area
    rect(slide, 0.6, 0.92, 12.13, 0.005, fill=RGBColor(0x20, 0x2B, 0x21))

def section_band(slide, label):
    # Right side: vertical kicker
    tb = slide.shapes.add_textbox(Inches(12.55), Inches(0.42), Inches(0.6), Inches(0.4))
    tf = tb.text_frame; tf.margin_left=tf.margin_right=0; tf.margin_top=tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = label.upper()
    r.font.name=FONT_B; r.font.size=Pt(9); r.font.bold=True; r.font.color.rgb=GOLD
    rPr = r._r.get_or_add_rPr(); rPr.set("spc","300")

# ---------- SLIDE 1 — CAPA ----------
def slide_cover():
    s = add_slide()
    # subtle vignette via darker rectangle bands
    rect(s, 0, 0, 13.333, 7.5, fill=BG)
    # soft radial-ish center using overlapping translucent ovals isn't trivial; instead use layered rects
    rect(s, 0, 5.9, 13.333, 1.6, fill=RGBColor(0x05, 0x0A, 0x06))

    # tiny top kicker with animated-dot feel (static dot)
    gold_dot(s, 0.62, 0.64, 0.09)
    text(s, 0.82, 0.5, 6, 0.35, "ESTRATÉGIA COMERCIAL B2B  ·  MENTORIA  ·  EXECUÇÃO",
         font=FONT_B, size=9, bold=True, color=GOLD, tracking=350)

    # brand wordmark top right
    text(s, 11.1, 0.45, 1.6, 0.4, "Outpace", font=FONT_H, size=20, bold=True, color=GOLD, align="right")

    # main title — centered mid
    title_serif(s, 0.9, 2.35, 11.5, 1.3,
                [[("Só crescer ", {}), ("não basta.", {"italic": True, "color": GOLD_L})]],
                size=60, align="center", line_spacing=1.0)
    title_serif(s, 0.9, 3.35, 11.5, 1.3,
                [[("É preciso ", {}), ("ultrapassar.", {"italic": True, "color": GOLD_L})]],
                size=60, align="center", line_spacing=1.0)

    # subline
    body(s, 2, 4.85, 9.3, 0.8,
         "Sistema comercial para indústrias e empresas B2B que querem parar de depender "
         "de esforço individual — e começar a operar com método, previsibilidade e padrão.",
         size=15, color=LIGHT, align="center", line_spacing=1.65)

    # footer meta on cover
    rect(s, 0.6, 6.85, 12.13, 0.005, fill=RGBColor(0x20, 0x2B, 0x21))
    text(s, 0.6, 7.0, 6, 0.3, "APRESENTAÇÃO  ·  SESSÃO ESTRATÉGICA",
         font=FONT_B, size=9, bold=True, color=MUTED, tracking=300)
    text(s, 7.3, 7.0, 5.4, 0.3, "CONFIDENCIAL  ·  MATERIAL DO CLOSER",
         font=FONT_B, size=9, bold=True, color=MUTED, tracking=300, align="right")

# ---------- SLIDE 2 — O MOMENTO ----------
def slide_momento(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 01 · O momento")

    kicker(s, 0.6, 1.3, "O cenário que ninguém discute")
    title_serif(s, 0.6, 1.75, 8.5, 2.2,
                [[("O comercial brasileiro\ncresce por "), ("esforço", {"italic": True, "color": GOLD_L}),
                  (", não por ",), ("método.", {"italic": True, "color": GOLD_L})]],
                size=46)

    body(s, 0.6, 4.1, 8.2, 2.4,
         "A maioria das empresas B2B e indústrias bateu em um teto invisível. "
         "Não é falta de talento. Não é falta de produto. É falta de um sistema "
         "comercial que sustente crescimento sem depender de quem é bom naquele dia.",
         size=15, line_spacing=1.7)

    # big stat block on the right
    rect(s, 9.35, 1.55, 3.4, 5.0, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
    rect(s, 9.35, 1.55, 0.06, 5.0, fill=GOLD)
    text(s, 9.7, 1.9, 3, 0.3, "A REALIDADE EM 3 NÚMEROS",
         font=FONT_B, size=9, bold=True, color=GOLD, tracking=250)

    stats = [
        ("72%", "das empresas B2B não têm previsibilidade de receita acima de 60 dias.",  2.55),
        ("3 em 5", "vendedores entregam resultado por dom pessoal, não por playbook.",   4.05),
        ("< 20%", "dos CRMs implantados são de fato usados como ferramenta de gestão.",  5.55),
    ]
    for v, lbl, yy in stats:
        text(s, 9.7, yy, 3, 0.7, v, font=FONT_H, size=32, bold=True, color=GOLD_L)
        body(s, 9.7, yy+0.72, 2.9, 0.8, lbl, size=11, color=LIGHT, line_spacing=1.45)

    footer(s, page, total)

# ---------- SLIDE 3 — SINTOMAS ----------
def slide_sintomas(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 02 · O diagnóstico")

    kicker(s, 0.6, 1.3, "Você provavelmente reconhece")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("Seis sintomas de um "), ("comercial que trava.", {"italic": True, "color": GOLD_L})]],
                size=36)

    sintomas = [
        ("01", "Pipeline sem controle",     "Ninguém sabe de verdade quanto vai entrar no próximo trimestre."),
        ("02", "CRM abandonado",            "Campos em branco, negociações perdidas em WhatsApp e planilha paralela."),
        ("03", "Venda por talento",         "Dois ou três vendedores carregam a meta — se saem, o resultado cai."),
        ("04", "Forecast é chute",          "Reunião de resultado vira reunião de explicação. Sem dado que confronte."),
        ("05", "Processo na cabeça",        "Playbook e padrão de abordagem não existem por escrito em lugar nenhum."),
        ("06", "Gestor bombeiro",           "Liderança apaga incêndio no dia a dia. Tempo zero para estratégia real."),
    ]
    x0, y0, cw, ch, gx, gy = 0.6, 3.55, 3.96, 1.5, 0.2, 0.2
    for i, (n, h, d) in enumerate(sintomas):
        col = i % 3; row = i // 3
        x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
        rect(s, x, y, cw, ch, fill=BG2, line=RGBColor(0x24, 0x30, 0x1E))
        text(s, x+0.3, y+0.22, 0.8, 0.5, n, font=FONT_H, size=20, bold=True, color=GOLD)
        text(s, x+1.0, y+0.28, cw-1.2, 0.45, h, font=FONT_B, size=13, bold=True, color=CREAM)
        body(s, x+1.0, y+0.72, cw-1.2, 0.7, d, size=10.5, color=LIGHT, line_spacing=1.5)

    footer(s, page, total)

# ---------- SLIDE 4 — A VIRADA ----------
def slide_virada(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 03 · A virada")

    kicker(s, 0.6, 1.35, "A tese que orienta tudo que fazemos")

    # Big pull quote style
    text(s, 0.6, 2.0, 0.6, 1.2, "“", font=FONT_H, size=110, color=GOLD, line_spacing=0.9)

    title_serif(s, 1.3, 2.1, 11.3, 3.2,
                [[("Comercial não é heroísmo.\nComercial é "), ("sistema.", {"italic": True, "color": GOLD_L})]],
                size=58, line_spacing=1.05)

    body(s, 1.3, 5.15, 10.8, 1.2,
         "Quando o comercial vira sistema, o resultado para de depender de quem acordou "
         "inspirado. Começa a depender de processo, indicador e cadência — os três lugares "
         "onde a Outpace entra para operar junto com você.",
         size=15, color=LIGHT, line_spacing=1.7)

    text(s, 1.3, 6.55, 5, 0.3, "— JEAN MACHADO  ·  FUNDADOR, OUTPACE",
         font=FONT_B, size=9, bold=True, color=GOLD, tracking=300)

    footer(s, page, total)

# ---------- SLIDE 5 — QUEM CONDUZ ----------
def slide_quem(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 04 · Quem conduz")

    # Left: photo block
    rect(s, 0.6, 1.3, 4.8, 5.45, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
    photo = os.path.join(IMG, "foto-perfil.png")
    if os.path.exists(photo):
        s.shapes.add_picture(photo, Inches(0.8), Inches(1.5), width=Inches(4.4), height=Inches(4.6))
    # badge
    rect(s, 0.8, 6.15, 4.4, 0.5, fill=RGBColor(0x14, 0x22, 0x16), line=GOLD, line_w=0.5)
    text(s, 1.0, 6.22, 3, 0.2, "JEAN MACHADO", font=FONT_B, size=10, bold=True, color=CREAM, tracking=200)
    text(s, 1.0, 6.42, 4.2, 0.2, "Fundador · Estrategista Comercial B2B",
         font=FONT_B, size=9, color=LIGHT)

    # Right: content
    kicker(s, 5.9, 1.3, "De promotor a mentor")
    title_serif(s, 5.9, 1.75, 7, 1.6,
                [[("O estrategista por trás\nda "), ("aceleração.", {"italic": True, "color": GOLD_L})]],
                size=34)

    body(s, 5.9, 3.5, 6.9, 3.0,
         "18 anos estruturando operações comerciais em indústrias e B2B. Da porta do cliente à "
         "direção comercial — conheço cada camada do problema e cada ponto onde a operação trava.\n\n"
         "Atuo dentro do negócio, não em cima dele. Organizo o processo, construo o playbook e "
         "acompanho a execução com o time até o resultado aparecer de forma consistente.",
         size=13, line_spacing=1.7)

    # chips
    tags = ["Indústria B2B", "Playbook Comercial", "Pipeline & CRM", "Gestão de Times", "18 anos"]
    cx = 5.9
    for t in tags:
        w = 0.22 + 0.09*len(t)
        rect(s, cx, 6.25, w, 0.36, fill=BG2, line=GOLD_D, line_w=0.5)
        text(s, cx, 6.25, w, 0.36, t, font=FONT_B, size=9, color=CREAM, align="center", valign="middle")
        cx += w + 0.12

    footer(s, page, total)

# ---------- SLIDE 6 — LOGOS ----------
def slide_logos(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 05 · Quem já ultrapassou")

    kicker(s, 0.6, 1.3, "Empresas que escolheram método")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("De multinacionais a indústrias —\nresultado "), ("em operação real.", {"italic": True, "color": GOLD_L})]],
                size=34)

    # logo band: 8 logos on white cards for contrast
    logos = [
        ("logo-claro.png",              "Claro"),
        ("logo-vivo.png",               "Vivo"),
        ("logo-lebes.png",              "Lebes"),
        ("logo-strazmaq.png",           "Strazmaq"),
        ("logo-lepetit.png",            "Le Petit"),
        ("logo-villerose.png",          "Ville Rose"),
        ("logo-ibs.png",                "IBS"),
        ("logo-escola-empresarios.png", "Escola de Empresários"),
    ]
    x0, y0 = 0.6, 3.9
    cw, ch, gx = 2.95, 1.35, 0.12
    for i, (fn, name) in enumerate(logos):
        col = i % 4; row = i // 4
        x = x0 + col*(cw+gx); y = y0 + row*(ch + 0.25)
        rect(s, x, y, cw, ch, fill=WHITE, line=RGBColor(0x2A, 0x3A, 0x24))
        p = os.path.join(IMG, fn)
        if os.path.exists(p):
            # center-fit inside card with padding
            s.shapes.add_picture(p, Inches(x+0.35), Inches(y+0.22),
                                 width=Inches(cw-0.7), height=Inches(ch-0.44))

    body(s, 0.6, 6.75, 12, 0.4,
         "Indústrias, telecom, varejo e serviços — a tese se repete em qualquer B2B sério.",
         size=11, color=MUTED, align="center")

    footer(s, page, total)

# ---------- SLIDE 7 — 3 FRENTES ----------
def slide_frentes(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 06 · A arquitetura")

    kicker(s, 0.6, 1.3, "Três frentes, uma máquina comercial completa")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("Como a Outpace "), ("entra na sua operação.", {"italic": True, "color": GOLD_L})]],
                size=34)

    produtos = [
        ("Partner Growth", "PARCERIA",
         "Presença estratégica, tática e operacional. Do CEO ao vendedor — "
         "fechamos o gap entre decisão e execução.",
         ["Reunião mensal com o CEO", "Semanal com líderes", "1:1 com vendedores",
          "Diagnóstico + Playbook", "Acompanhamento contínuo"]),
        ("SellAs Growth", "TIME COMERCIAL",
         "Foco no time. Reuniões quinzenais, playbook de vendas e assinatura "
         "contínua para melhorar execução e conversão.",
         ["Quinzenal com o time", "Quinzenal com gestores", "Playbook de vendas",
          "Padrão de abordagem", "Diagnóstico de execução"]),
        ("VendeAi", "CRM & OPERAÇÃO",
         "Implantação de CRM, ajuste de funil, automações e acompanhamento — "
         "para o CRM virar gestão, não cadastro.",
         ["Implantação de CRM", "Funil e etapas", "Automações", "Playbook no CRM",
          "Organização de dados"]),
    ]
    x0, y0, cw, ch, gx = 0.6, 3.4, 4.05, 3.55, 0.12
    for i, (nome, kk, desc, feats) in enumerate(produtos):
        x = x0 + i*(cw+gx); y = y0
        rect(s, x, y, cw, ch, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
        rect(s, x, y, 0.06, ch, fill=GOLD)  # left rail
        text(s, x+0.35, y+0.28, cw-0.5, 0.3, kk, font=FONT_B, size=9, bold=True,
             color=GOLD, tracking=300)
        text(s, x+0.35, y+0.58, cw-0.5, 0.55, nome, font=FONT_H, size=22, bold=True, color=CREAM)
        body(s, x+0.35, y+1.2, cw-0.6, 1.0, desc, size=11.5, color=LIGHT, line_spacing=1.55)
        # feats
        fy = y + 2.25
        for f in feats:
            gold_dot(s, x+0.38, fy+0.12, 0.06)
            text(s, x+0.58, fy, cw-0.8, 0.22, f, font=FONT_B, size=10.5, color=CREAM)
            fy += 0.25

    footer(s, page, total)

# ---------- PRODUCT DETAIL FACTORY ----------
def slide_produto(page, total, *, capnum, kk, nome, tagline, promessa, entregaveis, resultados, chapter_label):
    s = add_slide()
    header_mark(s); section_band(s, chapter_label)

    # Left: big name
    kicker(s, 0.6, 1.3, kk)
    title_serif(s, 0.6, 1.7, 7, 2.2,
                [[(nome + "\n"), (tagline, {"italic": True, "color": GOLD_L, "size": 38})]],
                size=52, line_spacing=1.0)
    body(s, 0.6, 4.05, 6.8, 1.4, promessa, size=13.5, color=LIGHT, line_spacing=1.7)

    # big product seal on left
    rect(s, 0.6, 5.7, 1.1, 1.1, fill=BG2, line=GOLD)
    text(s, 0.6, 5.7, 1.1, 1.1, str(capnum), font=FONT_H, size=40, bold=True,
         color=GOLD, align="center", valign="middle")
    text(s, 1.85, 5.92, 4.5, 0.3, "PRODUTO OUTPACE", font=FONT_B, size=9,
         bold=True, color=GOLD, tracking=300)
    text(s, 1.85, 6.22, 4.5, 0.3, nome.upper(), font=FONT_B, size=11,
         bold=True, color=CREAM, tracking=200)
    text(s, 1.85, 6.48, 4.5, 0.3, "Modalidade por assinatura",
         font=FONT_B, size=10, color=LIGHT)

    # Right: deliverables list
    rect(s, 7.7, 1.3, 5.05, 5.6, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
    text(s, 7.95, 1.55, 4.5, 0.3, "O QUE ESTÁ INCLUÍDO", font=FONT_B, size=9,
         bold=True, color=GOLD, tracking=250)
    text(s, 7.95, 1.85, 4.7, 0.45, "Entregáveis principais",
         font=FONT_H, size=18, bold=True, color=CREAM)

    yy = 2.55
    for i, (h, d) in enumerate(entregaveis, 1):
        text(s, 7.95, yy, 0.5, 0.3, f"{i:02d}", font=FONT_H, size=14, bold=True, color=GOLD)
        text(s, 8.5, yy, 4.1, 0.3, h, font=FONT_B, size=12, bold=True, color=CREAM)
        body(s, 8.5, yy+0.28, 4.1, 0.5, d, size=10, color=LIGHT, line_spacing=1.45)
        yy += 0.72

    # Bottom strip: outcomes
    rect(s, 0.6, 7.02, 12.13, 0.005, fill=RGBColor(0x20, 0x2B, 0x21))
    bx = 0.6
    for lbl, desc in resultados:
        text(s, bx, 7.1, 1.3, 0.25, lbl.upper(), font=FONT_B, size=8.5, bold=True,
             color=GOLD, tracking=300)
        bx += 4.05

    footer(s, page, total)

# ---------- SLIDE 11 — MÉTODO ----------
def slide_metodo(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 10 · O método")

    kicker(s, 0.6, 1.3, "Como a Outpace trabalha")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("Do "), ("diagnóstico", {"italic": True, "color": GOLD_L}),
                  (" à "), ("transformação.", {"italic": True, "color": GOLD_L})]],
                size=38)

    etapas = [
        ("01", "Diagnóstico",   "Mapeamos time, processo, funil e cultura com dado.",  "Semana 1–2"),
        ("02", "Estratégia",    "Plano de aceleração customizado ao seu estágio.",      "Semana 3–4"),
        ("03", "Implementação", "Entramos em campo: rituais, 1:1s, playbook, cadência.", "Mês 2–3"),
        ("04", "Resultados",    "Medimos, ajustamos, sustentamos. Sem relatório de gaveta.", "Mês 4+"),
    ]
    x0, y0, cw, ch, gx = 0.6, 3.6, 2.95, 3.0, 0.15
    # connecting hairline
    rect(s, 0.9, 3.95, 11.6, 0.01, fill=RGBColor(0x35, 0x2D, 0x14))
    for i, (n, h, d, when) in enumerate(etapas):
        x = x0 + i*(cw+gx); y = y0
        # number circle
        oval(s, x+0.0, y+0.1, 0.75, 0.75, fill=BG2, line=GOLD)
        text(s, x, y+0.1, 0.75, 0.75, n, font=FONT_H, size=16, bold=True,
             color=GOLD, align="center", valign="middle")
        text(s, x+0.9, y+0.25, 2, 0.3, when, font=FONT_B, size=9, bold=True,
             color=MUTED, tracking=250)
        text(s, x, y+1.05, cw-0.05, 0.5, h, font=FONT_H, size=22, bold=True, color=CREAM)
        body(s, x, y+1.65, cw-0.1, 1.3, d, size=12, color=LIGHT, line_spacing=1.6)

    footer(s, page, total)

# ---------- SLIDE 12 — RESULTADOS ----------
def slide_resultados(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 11 · Prova em números")

    kicker(s, 0.6, 1.3, "O que acontece quando o comercial vira sistema")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("Resultado "), ("tangível.", {"italic": True, "color": GOLD_L}),
                  (" Mensurável."), (" Sustentável.", {"italic": True, "color": GOLD_L})]],
                size=36)

    stats = [
        ("+30%",    "de previsibilidade de receita instalada em até 90 dias."),
        ("40%",     "de crescimento em um semestre — caso Lima Equipamentos."),
        ("3 níveis","de presença ativa: CEO, gestão e operação, em cadência semanal."),
        ("18 anos", "de experiência aplicada dentro do seu negócio, não em PPT."),
    ]
    x0, y0, cw, ch, gx = 0.6, 3.6, 2.95, 3.0, 0.15
    for i, (v, d) in enumerate(stats):
        x = x0 + i*(cw+gx); y = y0
        rect(s, x, y, cw, ch, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
        rect(s, x, y, 0.06, ch, fill=GOLD)
        text(s, x+0.3, y+0.5, cw-0.4, 1.3, v, font=FONT_H, size=48, bold=True,
             color=GOLD_L, line_spacing=1.0)
        rect(s, x+0.3, y+1.85, 0.4, 0.02, fill=GOLD)
        body(s, x+0.3, y+2.05, cw-0.5, 0.85, d, size=11.5, color=CREAM, line_spacing=1.6)

    footer(s, page, total)

# ---------- SLIDE 13 — TESTIMONIAL ----------
def slide_prova(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 12 · Prova real")

    kicker(s, 0.6, 1.3, "Um cliente. Uma operação. Um resultado.")

    # big quote card
    rect(s, 0.6, 2.0, 12.13, 4.65, fill=BG2, line=RGBColor(0x2A, 0x3A, 0x24))
    text(s, 0.95, 2.1, 1, 1.8, "“", font=FONT_H, size=150, color=GOLD, line_spacing=0.9)

    title_serif(s, 2.0, 2.55, 10.5, 2.6,
                "“Nosso time evoluiu muito, nossa coordenadora tem outra mentalidade, "
                "e temos novos canais de aquisição.”",
                size=26, line_spacing=1.35)

    # signature + logo card
    lp = os.path.join(IMG, "logo-strazmaq.png")
    rect(s, 2.0, 5.4, 1.8, 0.9, fill=WHITE, line=RGBColor(0x2A, 0x3A, 0x24))
    if os.path.exists(lp):
        s.shapes.add_picture(lp, Inches(2.1), Inches(5.5), width=Inches(1.6), height=Inches(0.7))
    text(s, 4.05, 5.5, 5, 0.35, "ALCINEI",
         font=FONT_B, size=12, bold=True, color=CREAM, tracking=150)
    text(s, 4.05, 5.85, 7, 0.35, "CEO & Diretor  ·  Strazmaq  ·  Indústria",
         font=FONT_B, size=11, color=LIGHT)
    text(s, 4.05, 6.1, 7, 0.25, "★★★★★",
         font=FONT_B, size=11, color=GOLD)

    footer(s, page, total)

# ---------- SLIDE 14 — POR QUE OUTPACE ----------
def slide_porque(page, total):
    s = add_slide()
    header_mark(s); section_band(s, "Capítulo 13 · Diferencial")

    kicker(s, 0.6, 1.3, "O que separa a Outpace de consultoria tradicional")
    title_serif(s, 0.6, 1.75, 12, 1.4,
                [[("Não entregamos "), ("relatório.", {"italic": True, "color": GOLD_L}),
                  (" Entregamos "), ("operação.", {"italic": True, "color": GOLD_L})]],
                size=36)

    diffs = [
        ("Dentro da operação",
         "Trabalhamos no seu negócio, não sobre ele. Reuniões, 1:1s, rituais e campo."),
        ("Três níveis ao mesmo tempo",
         "CEO, gestão e vendedor na mesma cadência. Estratégia e execução sincronizadas."),
        ("Método específico B2B",
         "Indústria, serviços e canais — o playbook é ajustado à sua realidade, não copiado."),
        ("Continuidade por assinatura",
         "Não somos projeto de 90 dias. Somos parceiro contínuo até virar hábito do time."),
        ("Dados acima de achismo",
         "Pipeline, KPIs e forecast com rastreabilidade. Decisão com base em número real."),
        ("Acervo de empresas sérias",
         "Claro, Vivo, Strazmaq, Lebes — a tese funciona de multinacional a indústria média."),
    ]
    x0, y0, cw, ch, gx, gy = 0.6, 3.45, 4.05, 1.65, 0.12, 0.15
    for i, (h, d) in enumerate(diffs):
        col = i % 3; row = i // 3
        x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
        rect(s, x, y, cw, ch, fill=BG2, line=RGBColor(0x24, 0x30, 0x1E))
        rect(s, x, y, cw, 0.04, fill=GOLD)
        text(s, x+0.3, y+0.25, cw-0.5, 0.4, h, font=FONT_H, size=15, bold=True, color=CREAM)
        body(s, x+0.3, y+0.75, cw-0.5, 0.8, d, size=11, color=LIGHT, line_spacing=1.55)

    footer(s, page, total)

# ---------- SLIDE 15 — PRÓXIMO PASSO ----------
def slide_cta(page, total):
    s = add_slide()
    # solid deep bg
    rect(s, 0, 0, 13.333, 7.5, fill=RGBColor(0x05, 0x0A, 0x06))
    # header mark
    text(s, 0.6, 0.42, 2, 0.4, "Outpace", font=FONT_H, size=18, bold=True, color=GOLD)
    rect(s, 0.6, 0.92, 12.13, 0.005, fill=RGBColor(0x20, 0x2B, 0x21))
    section_band(s, "Próximo passo")

    kicker(s, 0.6, 1.9, "Vamos desenhar seu plano de aceleração", size=11)

    title_serif(s, 0.6, 2.3, 12, 2.5,
                [[("O próximo trimestre\ncomeça com uma "),
                  ("decisão.", {"italic": True, "color": GOLD_L})]],
                size=58, line_spacing=1.02)

    body(s, 0.6, 4.7, 11, 1.0,
         "Em 45 minutos de sessão estratégica, mapeamos os três gargalos da sua operação "
         "comercial — e desenhamos o plano de como a Outpace entraria para resolver.",
         size=15, color=LIGHT, line_spacing=1.65)

    # Two-column next steps
    rect(s, 0.6, 5.9, 5.9, 1.1, fill=BG2, line=GOLD)
    rect(s, 0.6, 5.9, 0.06, 1.1, fill=GOLD)
    text(s, 0.85, 6.05, 5, 0.28, "1. AGENDAMOS HOJE",
         font=FONT_B, size=10, bold=True, color=GOLD, tracking=250)
    text(s, 0.85, 6.35, 5.6, 0.5, "Reserve a sessão antes do final desta conversa.",
         font=FONT_B, size=13, bold=True, color=CREAM)
    text(s, 0.85, 6.65, 5.6, 0.3, "Closer apresenta agenda em tela.",
         font=FONT_B, size=11, color=LIGHT)

    rect(s, 6.83, 5.9, 5.9, 1.1, fill=BG2, line=RGBColor(0x3A, 0x36, 0x22))
    rect(s, 6.83, 5.9, 0.06, 1.1, fill=GOLD_D)
    text(s, 7.08, 6.05, 5, 0.28, "2. OU NOS VOLTAMOS EM 24H",
         font=FONT_B, size=10, bold=True, color=GOLD, tracking=250)
    text(s, 7.08, 6.35, 5.6, 0.5, "Te mando um resumo executivo da sessão.",
         font=FONT_B, size=13, bold=True, color=CREAM)
    text(s, 7.08, 6.65, 5.6, 0.3, "+ proposta comercial dos produtos escolhidos.",
         font=FONT_B, size=11, color=LIGHT)

    text(s, 0.6, 7.12, 12.13, 0.3,
         "OUTPACE  ·  SÓ CRESCER NÃO BASTA. É PRECISO ULTRAPASSAR.",
         font=FONT_B, size=9, bold=True, color=GOLD, tracking=350, align="center")

# ---------- BUILD ----------
TOTAL = 15
slide_cover()                                                                    # 01
slide_momento(2, TOTAL)                                                           # 02
slide_sintomas(3, TOTAL)                                                          # 03
slide_virada(4, TOTAL)                                                            # 04
slide_quem(5, TOTAL)                                                              # 05
slide_logos(6, TOTAL)                                                             # 06
slide_frentes(7, TOTAL)                                                           # 07

slide_produto(8, TOTAL, capnum=1, kk="PARCERIA ESTRATÉGICA",
    nome="Partner Growth", tagline="que entra na operação.",
    promessa="Parceria estratégica e operacional em três níveis — CEO, gestores e time — "
             "para que a estratégia vire execução e a execução vire resultado previsível.",
    entregaveis=[
        ("Reunião mensal com o CEO",   "Metas, canais e decisões. O rumo sempre de pé."),
        ("Semanal com líderes",        "KPIs e pipeline analisados em cadência fixa."),
        ("1:1 com vendedores + PDI",   "Acompanhamento individual com plano próprio."),
        ("Diagnóstico da operação",    "Funil, dado e gargalo real — sem achismo."),
        ("Playbook e processos",       "Scripts, objeção, follow-up e qualificação documentados."),
        ("Acompanhamento contínuo",    "Presença semanal, dentro da operação."),
    ],
    resultados=[("+30% previsibilidade","a"),("3 níveis de atuação","b"),("Cadência semanal","c")],
    chapter_label="Produto 01 · Partner Growth")

slide_produto(9, TOTAL, capnum=2, kk="TIME COMERCIAL",
    nome="SellAs Growth", tagline="para o time vender com método.",
    promessa="Foco no time comercial. Reuniões quinzenais, playbook ajustado ao seu mercado "
             "e acompanhamento por assinatura para melhorar execução e conversão de verdade.",
    entregaveis=[
        ("Quinzenal com o time",       "Carteira, oportunidades e próximas ações."),
        ("Quinzenal com gestores",     "Pipeline, KPIs e decisões táticas."),
        ("Playbook de vendas",         "Scripts, cadências e objeções, prático e ajustado."),
        ("Padrão de abordagem",        "Processo definido — menos dependência de talento."),
        ("Acompanhamento contínuo",    "Presença mensal fixa, não pontual."),
        ("Diagnóstico de execução",    "Identificamos a quebra: prospecção, proposta ou fechamento."),
    ],
    resultados=[("Conversão real","a"),("Padrão instalado","b"),("Gestão com dado","c")],
    chapter_label="Produto 02 · SellAs Growth")

slide_produto(10, TOTAL, capnum=3, kk="CRM & OPERAÇÃO",
    nome="VendeAi", tagline="CRM virando gestão de verdade.",
    promessa="Implantação de CRM, ajuste de funil, automações e playbook dentro da ferramenta — "
             "mais acompanhamento contínuo para o CRM virar sistema de gestão, não cadastro.",
    entregaveis=[
        ("Implantação de CRM",         "Configuração completa: funil, etapas, campos e regras."),
        ("Ajustes de funil",           "Cada etapa com critério de saída e responsável."),
        ("Automações",                 "Follow-up, lembretes, roteamento e tarefas automáticas."),
        ("Playbook por área",          "Processo documentado dentro do CRM, não na cabeça de alguém."),
        ("Organização da operação",    "Dados limpos, pipeline saudável, sem planilha paralela."),
        ("Acompanhamento contínuo",    "Por assinatura, para o CRM manter padrão."),
    ],
    resultados=[("CRM em uso real","a"),("Pipeline visível","b"),("Menos retrabalho","c")],
    chapter_label="Produto 03 · VendeAi")

slide_metodo(11, TOTAL)     # 11
slide_resultados(12, TOTAL) # 12
slide_prova(13, TOTAL)      # 13
slide_porque(14, TOTAL)     # 14
slide_cta(15, TOTAL)        # 15

out = "/Users/jeanmachadopazdasilva/sellas-website/apresentacao/Outpace-Sessao-Estrategica.pptx"
prs.save(out)
print("SAVED", out)
